"""
Deliverable generation — turn an answer into a downloadable, attractive file.

Builds a **DOCX** (report/notes) or **PPTX** (slide deck) from a simple structured spec, with
**charts** (matplotlib) and **images** embedded so the output is easy to read and looks good —
for when a text card is not the best way to explain something (a study summary, a plan, a
comparison, a report).

Design:
- Pure builders take a `spec` dict and return file BYTES; storage/serving is separate.
- ALL heavy libraries (python-docx, python-pptx, matplotlib, Pillow) are imported LAZILY inside
  the functions, so importing this module never fails and the app runs even if they are not
  installed yet. `libraries_available()` reports what can be produced; callers degrade to a
  friendly message instead of crashing.

Spec shape (both formats share it):
  {
    "title": str, "subtitle": str, "author": "Nipun.AI",
    "sections": [
      {"heading": str,
       "bullets": [str], "paragraphs": [str],
       "chart": {"type": "bar|line|pie", "title": str, "labels": [str], "values": [num],
                 "series_label": str} | None,
       "image_bytes": bytes | None}          # embedded picture (already fetched)
    ]
  }
"""

from __future__ import annotations

import io

import structlog

log = structlog.get_logger("synthesis.filegen")


def _has(mod: str) -> bool:
    import importlib.util
    return importlib.util.find_spec(mod) is not None


def libraries_available() -> dict:
    """Which output formats can actually be produced right now."""
    docx = _has("docx")
    pptx = _has("pptx")
    charts = _has("matplotlib")
    return {"docx": docx, "pptx": pptx, "charts": charts,
            "any": docx or pptx}


# ── Charts ────────────────────────────────────────────────────────────────────

def render_chart_png(chart: dict) -> bytes | None:
    """Render a bar/line/pie chart to PNG bytes (matplotlib). None on any failure."""
    if not chart or not _has("matplotlib"):
        return None
    labels = [str(x) for x in (chart.get("labels") or [])]
    values = [float(v) for v in (chart.get("values") or []) if _is_num(v)]
    if not values or len(values) != len(labels):
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")   # headless — no display needed
        import matplotlib.pyplot as plt

        kind = (chart.get("type") or "bar").lower()
        fig, ax = plt.subplots(figsize=(6.2, 3.6), dpi=150)
        color = "#C2703D"  # warm accent to match the product
        if kind == "pie":
            ax.pie(values, labels=labels, autopct="%1.0f%%", startangle=90)
            ax.axis("equal")
        elif kind == "line":
            ax.plot(labels, values, marker="o", color=color, linewidth=2)
            ax.grid(True, alpha=0.25)
        else:  # bar
            ax.bar(labels, values, color=color)
            ax.grid(True, axis="y", alpha=0.25)
        if chart.get("title"):
            ax.set_title(chart["title"], fontsize=12, fontweight="bold")
        if kind != "pie" and chart.get("series_label"):
            ax.set_ylabel(chart["series_label"])
        if kind != "pie" and len(max(labels, key=len, default="")) > 8:
            plt.xticks(rotation=30, ha="right", fontsize=8)
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()
    except Exception as exc:
        log.debug("chart_render_failed", error=str(exc))
        return None


def _is_num(v) -> bool:
    try:
        float(v)
        return True
    except (TypeError, ValueError):
        return False


def render_flow_png(steps: list[str], title: str = "") -> bytes | None:
    """Render a top-to-bottom flow of labelled boxes with arrows (a simple process diagram).
    Great for 'how X works' answers whose content is a sequence of steps. None if matplotlib
    is missing or there are fewer than 2 steps."""
    steps = [str(s).strip() for s in (steps or []) if str(s).strip()][:6]
    if len(steps) < 2 or not _has("matplotlib"):
        return None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.patches import FancyBboxPatch

        n = len(steps)
        fig, ax = plt.subplots(figsize=(5.4, 0.95 * n + (0.5 if title else 0.1)), dpi=150)
        ax.axis("off")
        ax.set_xlim(0, 1)
        ax.set_ylim(0, n)
        accent, fill = "#C2703D", "#FBEEE4"
        for i, label in enumerate(steps):
            cy = n - i - 0.5                      # top step highest
            ax.add_patch(FancyBboxPatch(
                (0.08, cy - 0.32), 0.84, 0.64,
                boxstyle="round,pad=0.01,rounding_size=0.08",
                linewidth=1.6, edgecolor=accent, facecolor=fill))
            ax.text(0.5, cy, label, ha="center", va="center", fontsize=10, wrap=True,
                    color="#3a2a20")
            if i < n - 1:
                ax.annotate("", xy=(0.5, cy - 0.34), xytext=(0.5, cy - 0.66),
                            arrowprops=dict(arrowstyle="-|>", color=accent, lw=1.8))
        if title:
            ax.set_title(title, fontsize=12, fontweight="bold", color="#3a2a20")
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight")
        plt.close(fig)
        return buf.getvalue()
    except Exception as exc:
        log.debug("flow_render_failed", error=str(exc))
        return None


# ── DOCX ──────────────────────────────────────────────────────────────────────

def build_docx(spec: dict) -> bytes | None:
    """Build a Word document from the spec. None if python-docx is unavailable."""
    if not _has("docx"):
        return None
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor

        doc = Document()
        title = doc.add_heading(spec.get("title", "Document"), level=0)
        if spec.get("subtitle"):
            sub = doc.add_paragraph(spec["subtitle"])
            sub.runs[0].italic = True
        for sec in spec.get("sections", []):
            if sec.get("heading"):
                doc.add_heading(sec["heading"], level=1)
            # The plain-language explanation of this section, first, as normal text.
            if sec.get("notes"):
                doc.add_paragraph(str(sec["notes"]))
            for para in sec.get("paragraphs", []) or []:
                doc.add_paragraph(str(para))
            for b in sec.get("bullets", []) or []:
                doc.add_paragraph(str(b), style="List Bullet")
            png = render_chart_png(sec.get("chart")) if sec.get("chart") else None
            if png:
                doc.add_picture(io.BytesIO(png), width=Inches(5.5))
            if sec.get("image_bytes"):
                try:
                    doc.add_picture(io.BytesIO(sec["image_bytes"]), width=Inches(4.5))
                except Exception:
                    pass
        foot = doc.add_paragraph(f"Generated by {spec.get('author', 'Nipun.AI')}")
        foot.runs[0].font.size = Pt(8)
        foot.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except Exception as exc:
        log.warning("docx_build_failed", error=str(exc))
        return None


# ── PPTX ──────────────────────────────────────────────────────────────────────

def build_pptx(spec: dict) -> bytes | None:
    """Build a slide deck from the spec (one content slide per section). None if
    python-pptx is unavailable."""
    if not _has("pptx"):
        return None
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = spec.get("title", "Presentation")
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = spec.get("subtitle", "") or spec.get("author", "Nipun.AI")

        blank = prs.slide_layouts[6]     # blank layout for full control
        content = prs.slide_layouts[1]   # title + content

        def _add_notes(slide, text):
            if text:
                try:
                    slide.notes_slide.notes_text_frame.text = str(text)
                except Exception:
                    pass

        for sec in spec.get("sections", []):
            png = render_chart_png(sec.get("chart")) if sec.get("chart") else None
            img = sec.get("image_bytes")
            bullets = [str(b) for b in (sec.get("bullets") or [])] + \
                      [str(p) for p in (sec.get("paragraphs") or [])]
            if png or img:
                s = prs.slides.add_slide(blank)
                tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                tf = tb.text_frame
                tf.text = sec.get("heading", "")
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True
                # bullets on the left, visual on the right
                if bullets:
                    body = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.6), Inches(5))
                    bf = body.text_frame
                    bf.word_wrap = True
                    for i, b in enumerate(bullets[:8]):
                        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
                        p.text = "• " + b
                        p.font.size = Pt(16)
                try:
                    s.shapes.add_picture(io.BytesIO(png or img), Inches(5.3), Inches(1.4), width=Inches(4.2))
                except Exception:
                    pass
            else:
                s = prs.slides.add_slide(content)
                s.shapes.title.text = sec.get("heading", "")
                body = s.placeholders[1].text_frame if len(s.placeholders) > 1 else None
                if body is not None:
                    for i, b in enumerate(bullets[:10]):
                        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                        p.text = b
                        p.font.size = Pt(18)
            # Attach the spoken explanation for this slide as speaker notes.
            _add_notes(s, sec.get("notes"))
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except Exception as exc:
        log.warning("pptx_build_failed", error=str(exc))
        return None


FORMATS = {
    "docx": ("build_docx",
             "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "pptx": ("build_pptx",
             "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
}


def build(fmt: str, spec: dict) -> tuple[bytes, str] | None:
    """Build `fmt` ('docx'|'pptx') from spec → (bytes, mime), or None if unavailable."""
    fmt = (fmt or "").lower()
    if fmt not in FORMATS:
        return None
    builder_name, mime = FORMATS[fmt]
    data = globals()[builder_name](spec)
    return (data, mime) if data else None
